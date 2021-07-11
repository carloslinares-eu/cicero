import cicero.cfg as cfg
import cicero.csvhandler as csv_handler
import pptx
import tkinter.messagebox
from pptx.enum.shapes import MSO_SHAPE_TYPE


def open_powerpoint_file(path_to_input):
    try:
        active_presentation = pptx.Presentation(path_to_input)
        return active_presentation
    except pptx.exc.PackageNotFoundError:
        refresh_error_message = "Failed to open the Powerpoint file"
        tkinter.messagebox.showerror(title="Critical error", message=refresh_error_message)
        raise IOError


def extract_text_from_powerpoint_and_write_csv(active_presentation):
    cfg.element_id = 0
    for cfg.current_slide in active_presentation.slides:
        extract_text_from_notes(cfg.current_slide)
        extract_text_from_simple_shapes(cfg.current_slide)
        extract_text_from_complex_shapes(cfg.current_slide)
        extract_text_from_tables(cfg.current_slide)


def extract_text_from_notes(slide):
    extract_paragraphs_from_shapes(slide, slide.notes_slide.notes_text_frame)


def extract_text_from_simple_shapes(slide):
    for shape in slide.shapes:
        if shape.has_text_frame:
            extract_paragraphs_from_shapes(slide, shape.text_frame)


def extract_text_from_complex_shapes(slide):
    for grouped_shape in slide.shapes:
        if grouped_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            extract_text_from_simple_shapes(slide)


def extract_text_from_tables(slide):
    for shape in slide.shapes:
        if shape.has_table:
            for cell in shape.table.iter_cells():
                extract_paragraphs_from_shapes(slide, cell)


def extract_paragraphs_from_shapes(slide, shape_text_frame):
    for paragraph in shape_text_frame.paragraphs:
        for run in paragraph.runs:
            inline_dictionary = {"slide_id": slide.slide_id, "element_id": cfg.element_id,
                                 "encoded": str(False), "text": run.text, "font_name": run.font.name,
                                 "font_size": run.font.size, "font_color": "N/A",
                                 "translation": "Pending"}
            csv_handler.write_line_to_repository(cfg.text_repository_file_path, inline_dictionary)
            cfg.element_id += 1


def replace_text_with_translation(active_presentation):
    cfg.element_id = 0
    for cfg.current_slide in active_presentation.slides:
        replace_notes_with_translation()
        replace_simple_shapes_with_translation()
        replace_complex_shapes_with_translation()
        replace_tables_with_translation()


def replace_notes_with_translation():
    replace_paragraphs_with_translation(cfg.current_slide.notes_slide.notes_text_frame)


def replace_simple_shapes_with_translation():
    for shape in cfg.current_slide.shapes:
        if shape.has_text_frame:
            replace_paragraphs_with_translation(shape.text_frame)


def replace_complex_shapes_with_translation():
    for grouped_shape in cfg.current_slide.shapes:
        if grouped_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            replace_simple_shapes_with_translation()


def replace_tables_with_translation():
    for shape in cfg.current_slide.shapes:
        if shape.has_table:
            for cell in shape.table.iter_cells():
                replace_paragraphs_with_translation(cell)


def replace_paragraphs_with_translation(shape_text_frame):
    for paragraph in shape_text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = cfg.repository_content[cfg.element_id]["translation"]
            if cfg.repository_content[cfg.element_id]["font_name"] != "":
                run.font.name = cfg.repository_content[cfg.element_id]["font_name"]
            if cfg.repository_content[cfg.element_id]["font_size"] != "":
                run.font.size = int(cfg.repository_content[cfg.element_id]["font_size"])
            cfg.element_id += 1
