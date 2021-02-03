import pptx


def extract_notes(pptx_file_path):
    active_pptx = pptx.Presentation(pptx_file_path)
    data = {"speaker_notes": []}

    for slide_number, slide in enumerate(active_pptx.slides):
        text_note = slide.notes_slide.notes_text_frame.text
        data["speaker_notes"].append({
            "slide_number": slide_number,
            "text": text_note
        })
    return data


def extract_all_text(pptx_file_path):
    active_pptx = pptx.Presentation(pptx_file_path)
    data = {"shapes_with_text": []}

    for slide_number, slide in enumerate(active_pptx.slides):
        for shape_id, shape in enumerate(slide.shapes):
            if hasattr(shape, "text"):
                data["shapes_with_text"].append({
                    "slide_number": slide_number,
                    "shape_id": shape_id,
                    "text": shape.text
                })
    return data


notes = extract_notes("/home/carlos/Documentos/cicero/200411 MBA AEROSP Sevilla Exercise 5 Channels (in class).pptx")
print(notes)

text = extract_all_text("/home/carlos/Documentos/cicero/200411 MBA AEROSP Sevilla Exercise 5 Channels (in class).pptx")
print(text)
