import tkinter
import tkinter.messagebox
import tkinter.scrolledtext
import tkinter.ttk
import tkinter.filedialog
import google.oauth2.service_account
import google.cloud.translate_v2
import google.auth.exceptions
import pptx


class TranslationApp:
    def __init__(self, icon_path, credentials=None):
        self.credentials = credentials
        self.icon = icon_path
        self.translate_client = google.cloud.translate_v2.Client(credentials=self.credentials)
        self.root = tkinter.Tk()
        self.root.title("CICERO - Universal Translator App")
        self.root.geometry("500x260")
        self.root.resizable(0, 0)
        self.root.iconbitmap(self.icon)
        self.large = ("Segoe UI", 10)
        self.medium = ("Segoe UI", 8)
        self.small = ("Segoe UI", 6)
        self.main_width = 60
        self.compatible_files = (("Powerpoint presentation", "*.pptx"), ("PDF file", "*.pdf"))
        self.path_to_input = None
        self.path_to_output = None
        self.file_type = None

        self.input_label = tkinter.Label(self.root, text="Input file", anchor="w", font=self.large)
        self.input_entry = tkinter.Entry(self.root, font=self.large, width=54)
        self.browse_button = tkinter.ttk.Button(self.root, text="Browse", command=self.set_input)

        self.origin_label = tkinter.ttk.Label(self.root, text="Origin Language", anchor="w", font=self.medium)
        self.origin_combo = tkinter.ttk.Combobox(self.root)
        self.origin_combo.configure(postcommand=self.origin_combo.configure(values=self.get_language_list()))

        self.target_label = tkinter.ttk.Label(self.root, text="Target Language", anchor="w", font=self.medium)
        self.target_combo = tkinter.ttk.Combobox(self.root)
        self.target_combo.configure(postcommand=self.target_combo.configure(values=self.get_language_list()))

        self.horizontal_separator = tkinter.ttk.Separator(self.root, orient="horizontal")

        self.output_label = tkinter.Label(self.root, text="Output file", anchor="w", font=self.large)
        self.output_entry = tkinter.Entry(self.root, font=self.large, width=54)
        self.save_button = tkinter.ttk.Button(self.root, text="Save", command=self.set_output)

        self.translate_button = tkinter.ttk.Button(self.root, text="Translate!", command=self.translate_powerpoint_file)

    def initialize(self):
        current_row = 0
        self.input_label.grid(padx=10, pady=10, row=current_row, column=0, columnspan=2, sticky="W")

        current_row += 1
        self.input_entry.grid(padx=10, pady=0, row=current_row, column=0, columnspan=2)
        self.browse_button.grid(padx=10, row=current_row, column=2, sticky="EW")

        current_row += 1
        self.origin_label.grid(padx=10, pady=10, row=current_row, column=0, sticky="W")
        self.origin_combo.grid(padx=10, row=current_row, column=1, columnspan=2, sticky="E", pady=5)

        current_row += 1
        self.horizontal_separator.grid(padx=10, row=current_row, column=0, columnspan=3, sticky="EW")

        current_row += 1
        self.output_label.grid(padx=10, pady=10, row=current_row, column=0, columnspan=2, sticky="W")

        current_row += 1
        self.output_entry.grid(padx=10, pady=0, row=current_row, column=0, columnspan=2)
        self.save_button.grid(padx=10, row=current_row, column=2, sticky="EW")

        current_row += 1
        self.target_label.grid(padx=10, pady=10, row=current_row, column=0, sticky="W")
        self.target_combo.grid(padx=10, row=current_row, column=1, columnspan=2, sticky="E", pady=5)

        current_row += 1
        self.translate_button.grid(padx=10, row=current_row, column=2, sticky="W", pady=5)

        self.root.mainloop()

    def set_input(self):
        prompt_title = "Open file to translate"
        self.path_to_input = tkinter.filedialog.askopenfile(filetypes=self.compatible_files, title=prompt_title)

        if self.path_to_input is not None:
            self.input_entry.delete(0, 'end')

        self.input_entry.insert(0, self.path_to_input.name)

        self.set_file_type()

    def set_file_type(self):
        self.file_type = None
        if self.path_to_input.name[-5:] == ".pptx":
            self.file_type = {("Powerpoint presentation", "*.pptx")}
        elif self.path_to_input.name[-4:] == ".pdf":
            self.file_type = {("PDF file", "*.pdf")}

    def set_output(self):
        prompt_title = "Save translated file to:"
        self.path_to_output = tkinter.filedialog.asksaveasfilename(filetypes=self.file_type, title=prompt_title)

        if self.path_to_input is not None:
            self.output_entry.delete(0, 'end')

        self.complete_output_path()
        self.output_entry.insert(0, self.path_to_output)

    def complete_output_path(self):
        if (self.file_type == {("Powerpoint presentation", "*.pptx")}) & (self.path_to_output[-5:] != ".pptx"):
            self.path_to_output += ".pptx"
        elif (self.file_type == {("PDF file", "*.pdf")}) & (self.path_to_output[-4:] != ".pdf"):
            self.path_to_output += ".pdf"

    def get_language_list(self, language_code=None):
        language_list = self.translate_client.get_languages(language_code)
        screen_name_list = []
        for language in language_list:
            screen_name_list.append(language.get("language") + " - " + language.get("name"))
        return screen_name_list

    def translate_powerpoint_file(self):
        source_file_path = self.input_entry.get()
        destination_file_path = self.output_entry.get()
        target_lang = self.target_combo.get()[0:2]
        active_presentation = None
        try:
            active_presentation = pptx.Presentation(source_file_path)
        except pptx.exc.PackageNotFoundError:
            refresh_error_message = "Failed to open the Powerpoint file"
            tkinter.messagebox.showerror(title="Critical error", message=refresh_error_message)

        for slide in active_presentation.slides:
            input_text = slide.notes_slide.notes_text_frame.text
            translation = self.translate_client.translate(input_text, target_language=target_lang)
            slide.notes_slide.notes_text_frame.text = translation["translatedText"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        translation = self.translate_client.translate(paragraph.text, target_language=target_lang)
                        paragraph.text = translation["translatedText"]
                elif shape.has_table:
                    for cell in shape.table.iter_cells():
                        for paragraph in cell.text_frame.paragraphs:
                            translation = self.translate_client.translate(paragraph.text, target_language=target_lang)
                            paragraph.text = translation["translatedText"]
                else:
                    continue

        active_presentation.save(destination_file_path)
        tkinter.messagebox.showinfo(title="Job finished", message="The file has been successfully translated")
